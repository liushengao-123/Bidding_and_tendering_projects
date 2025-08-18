import json
import os
import re
from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE_TYPE # 注意：此行在原代码中未被使用，可以安全移除

def extract_structured_text_from_pptx(file_path, start_page=None, end_page=None):
    """
    从具有特定模板的PowerPoint (.pptx)文件中精细化提取结构化文本。
    此版本能正确处理值跨越多行，并使用 '■' 作为新条目的明确分隔符。
    
    【优化】: start_page 和 end_page 参数是可选的。如果留空，则默认提取所有幻灯片。

    参数:
        file_path (str): PPTX文件的路径。
        start_page (int, optional): 开始提取的页码（从1开始）。如果为 None，则从第一页开始。默认为 None。
        end_page (int, optional): 结束提取的页码（包含此页）。如果为 None，则提取到最后一页。默认为 None。

    返回:
        dict: 一个字典，键是幻灯片标题，值是一个包含该标题下所有键值对的列表。
    """
    try:
        presentation = Presentation(file_path)
        structured_data = {}

        total_slides = len(presentation.slides)

        # --- 优化核心逻辑开始 ---
        # 确定最终要处理的幻灯片范围
        # 如果未提供 start_page，则从第一页 (1) 开始
        start = start_page if start_page is not None else 1
        # 如果未提供 end_page，则到最后一页 (total_slides) 结束
        end = end_page if end_page is not None else total_slides
        
        # 使用最终确定的 start 和 end 进行范围验证
        if start < 1 or end > total_slides or start > end:
            print(f"错误: 页码范围无效。您提供的范围是 {start_page}-{end_page}，但有效范围是 1 到 {total_slides}。")
            return {}
        # --- 优化核心逻辑结束 ---
            
        for slide_idx, slide in enumerate(presentation.slides):
            current_page_num = slide_idx + 1
            # 使用最终确定的 start 和 end 来过滤幻灯片
            if not (start <= current_page_num <= end):
                continue

            print(f"--- 正在处理第 {current_page_num} 页 (总共 {total_slides} 页) ---")
            
            title = f"未找到标题_第{current_page_num}页"
            text_containing_shapes = [
                s for s in slide.shapes 
                if s.has_text_frame and s.text_frame.text.strip()
            ]
            
            if not text_containing_shapes:
                print(f"警告: 第 {current_page_num} 页没有找到任何文本内容。")
                continue

            sorted_shapes = sorted(text_containing_shapes, key=lambda s: s.top)
            
            title_shape = sorted_shapes[0]
            if title_shape.text_frame.paragraphs:
                title = title_shape.text_frame.paragraphs[0].text.strip()
            else:
                 title = title_shape.text.strip().split('\n')[0]
            
            print(f"识别到标题: {title}")

            slide_content_list = []
            
            # 为了避免将标题也作为内容解析，从第二个形状开始处理
            content_shapes = sorted_shapes[1:] if len(sorted_shapes) > 1 else []

            for shape in content_shapes:
                all_lines = re.split(r'[\n\v\f\r]+', shape.text_frame.text)
                
                current_key = None
                current_value_lines = []

                def save_previous_item():
                    nonlocal current_key, current_value_lines
                    if current_key:
                        full_value = ' '.join(current_value_lines).strip()
                        if full_value:
                            slide_content_list.append({current_key: full_value})
                            print(f"  提取到 -> {{'{current_key}': '{full_value[:70]}...'}}")
                    current_key = None
                    current_value_lines = []

                for line in all_lines:
                    clean_line = line.strip()
                    if not clean_line:
                        continue

                    if clean_line.startswith('■'):
                        save_previous_item()
                        line_to_parse = clean_line.lstrip('■').strip()
                        parts = line_to_parse.replace('：', ':', 1).split(':', 1)
                        
                        if len(parts) == 2:
                            current_key = parts[0].strip()
                            current_value_lines.append(parts[1].strip())
                        else:
                            current_key = line_to_parse
                            current_value_lines = ['']
                    elif current_key:
                        current_value_lines.append(clean_line)
                
                save_previous_item()

            if title in structured_data:
                structured_data[title].extend(slide_content_list)
            else:
                structured_data[title] = slide_content_list

        return structured_data

    except Exception as e:
        import traceback
        print(f"处理PPT文件时发生严重错误: {e}")
        traceback.print_exc()
        return {}

# --- 使用示例 (已更新以演示新功能) ---
if __name__ == "__main__":
    # 请将路径替换为您的实际文件路径
    # pptx_file_path = 'E:/project/ppt精细化抽取/关于榆林分公司2025年-2026年员工食堂食材采购项目采购方案的汇报.pptx'
    pptx_file_path = 'E:\project\优化版本1\ppt精细化抽取\关于咸阳分公司2023-2024年度职工餐厅食材采购.pptx'
    #pptx_file_path = 'E:\project\ppt精细化抽取\关于咸阳分公司2023-2024年度职工餐厅食材采购.pptx'
    base_name = os.path.basename(pptx_file_path)
    file_name_without_ext = os.path.splitext(base_name)[0]
    
    if not os.path.exists(pptx_file_path):
        print(f"错误：输入文件不存在 -> '{pptx_file_path}'")
    else:
        # --- 示例1: 提取特定范围 (如第 3 到 5 页) ---
        # print("\n================== 示例1: 提取第 3-5 页 ==================")
        # output_json_path_partial = f"{file_name_without_ext}_extracted_pages_3-5.json"
        # extracted_data_partial = extract_structured_text_from_pptx(pptx_file_path, start_page=3, end_page=5)
        
        # if extracted_data_partial:
        #     try:
        #         with open(output_json_path_partial, 'w', encoding='utf-8') as f:
        #             json.dump(extracted_data_partial, f, indent=4, ensure_ascii=False)
        #         print(f"\n✅ 特定范围提取结果已成功保存到文件: {output_json_path_partial}")
        #     except IOError as e:
        #         print(f"\n❌ 保存文件时出错: {e}")
        # else:
        #     print("\n- 在指定范围内未提取到任何数据。")

        # --- 示例2: 提取全部 (不提供 start 和 end 参数) ---
        print("\n================== 示例2: 提取全部页面 ==================")
        output_json_path_all = f"{file_name_without_ext}_extracted_all.json"
        extracted_data_all = extract_structured_text_from_pptx(pptx_file_path) # 调用时不传 start/end
        
        if extracted_data_all:
            print("\n--- 最终提取结果 (全部页面预览) ---")
            # 为避免控制台输出过长，可以只打印键
            print("提取到的标题:", list(extracted_data_all.keys()))
            try:
                with open(output_json_path_all, 'w', encoding='utf-8') as f:
                    json.dump(extracted_data_all, f, indent=4, ensure_ascii=False)
                print(f"\n✅ 全部提取结果已成功保存到文件: {output_json_path_all}")
            except IOError as e:
                print(f"\n❌ 保存文件时出错: {e}")
        else:
            print("\n- 未提取到任何数据，不创建输出文件。")