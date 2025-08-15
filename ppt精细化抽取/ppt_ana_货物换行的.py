import json
import os
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_structured_text_from_pptx(file_path, start_page, end_page):
    """
    从具有特定模板的PowerPoint (.pptx)文件中精细化提取结构化文本。
    此版本能处理并行键值对（即一个段落内包含多个逻辑条目）。

    参数:
        file_path (str): PPTX文件的路径。
        start_page (int): 开始提取的页码（从1开始）。
        end_page (int): 结束提取的页码（包含此页）。

    返回:
        dict: 一个字典，键是幻灯片标题，值是一个包含该标题下所有键值对的列表。
    """
    try:
        presentation = Presentation(file_path)
        structured_data = {}

        total_slides = len(presentation.slides)
        if start_page < 1 or end_page > total_slides or start_page > end_page:
            print(f"错误: 页码范围无效。有效范围是 1 到 {total_slides}。")
            return {}
            
        for slide_idx, slide in enumerate(presentation.slides):
            current_page_num = slide_idx + 1
            if not (start_page <= current_page_num <= end_page):
                continue

            print(f"--- 正在处理第 {current_page_num} 页 ---")
            
            title = f"未找到标题_第{current_page_num}页"
            text_containing_shapes = [
                s for s in slide.shapes 
                if s.has_text_frame and s.text_frame.text.strip()
            ]
            
            if not text_containing_shapes:
                print(f"警告: 第 {current_page_num} 页没有找到任何文本内容。")
                continue

            sorted_shapes = sorted(text_containing_shapes, key=lambda s: s.top)
            
            title_shape = sorted_shapes[0]
            if title_shape.text_frame.paragraphs:
                title = title_shape.text_frame.paragraphs[0].text.strip()
            else:
                 title = title_shape.text.strip().split('\n')[0]
            
            print(f"识别到标题: {title}")

            slide_content_list = []
            
            # 【核心修正】重构解析逻辑以处理并行键值对
            for shape in sorted_shapes:
                for paragraph in shape.text_frame.paragraphs:
                    # 将一个段落的文本按多种换行符（\n, \u000b等）分割成多行
                    # re.split能够处理多种分隔符
                    lines = re.split(r'[\n\v\f\r]+', paragraph.text)
                    
                    # 临时变量存储多行value的情况
                    current_key = None
                    current_value_lines = []

                    def process_collected_item():
                        nonlocal current_key, current_value_lines
                        if current_key:
                            value = ' '.join(current_value_lines).strip()
                            if value:
                                slide_content_list.append({current_key: value})
                                print(f"  提取到 -> {{'{current_key}': '{value}'}}")
                        current_key = None
                        current_value_lines = []

                    for line in lines:
                        line = line.strip()
                        if not line or line == title:
                            continue

                        # 使用正则表达式匹配 "■ Key: Value" 或 "Key: Value" 格式
                        # 冒号可以是半角或全角
                        match = re.match(r'^[■\s]*(?P<key>[^：:]+)[：:](?P<value>.*)', line)
                        
                        if match:
                            # 如果找到一个新的key，先处理之前收集好的键值对
                            process_collected_item()
                            
                            # 开始收集新的键值对
                            current_key = match.group('key').strip()
                            current_value_lines.append(match.group('value').strip())
                        elif current_key:
                            # 如果这一行没有key，但我们正在收集一个value，就追加到value上
                            # 这处理了值跨越多行的情况
                            current_value_lines.append(line)
                    
                    # 处理循环结束时最后一个收集到的键值对
                    process_collected_item()

            if title in structured_data:
                structured_data[title].extend(slide_content_list)
            else:
                structured_data[title] = slide_content_list

        return structured_data

    except Exception as e:
        import traceback
        print(f"处理PPT文件时发生严重错误: {e}")
        traceback.print_exc()
        return {}

# --- 使用示例 ---
if __name__ == "__main__":
    # --- 配置区域 ---
    pptx_file_path = 'E:\project\ppt精细化抽取\关于榆林分公司2025年-2026年员工食堂食材采购项目采购方案的汇报.pptx'
    base_name = os.path.basename(pptx_file_path)
    file_name_without_ext = os.path.splitext(base_name)[0]
    output_json_path = f"{file_name_without_ext}_extracted.json"
    start = 3
    end = 5
    
    # --- 主逻辑 ---
    if not os.path.exists(pptx_file_path):
        print(f"错误：输入文件不存在 -> '{pptx_file_path}'")
        print("请将代码中的 'path/to/your/example.pptx' 替换为您的PPT文件的真实路径。")
    else:
        extracted_data = extract_structured_text_from_pptx(pptx_file_path, start, end)
        
        print("\n--- 最终提取结果 (控制台预览) ---")
        print(json.dumps(extracted_data, indent=4, ensure_ascii=False))

        # --- 文件保存逻辑 ---
        if extracted_data:
            try:
                with open(output_json_path, 'w', encoding='utf-8') as f:
                    json.dump(extracted_data, f, indent=4, ensure_ascii=False)
                print(f"\n✅ 提取结果已成功保存到文件: {output_json_path}")
            except IOError as e:
                print(f"\n❌ 保存文件时出错: {e}")
        else:
            print("\n- 未提取到任何数据，不创建输出文件。")