import json
import os
import re
from pptx import Presentation

def extract_structured_text_from_pptx(file_path, start_page, end_page):
    """
    从具有特定模板的PowerPoint (.pptx)文件中精细化提取结构化文本。
    V7 (决定版): 采用更精准的“肯定法”正则表达式，彻底解决次级键值对提取错误的问题。

    参数:
        file_path (str): PPTX文件的路径。
        start_page (int): 开始提取的页码（从1开始）。
        end_page (int): 结束提取的页码（包含此页）。

    返回:
        dict: 一个字典，键是幻灯片标题，值是一个包含该标题下所有键值对的列表。
    """
    try:
        presentation = Presentation(file_path)
        structured_data = {}

        total_slides = len(presentation.slides)
        if start_page < 1 or end_page > total_slides or start_page > end_page:
            print(f"错误: 页码范围无效。有效范围是 1 到 {total_slides}。")
            return {}
            
        for slide_idx, slide in enumerate(presentation.slides):
            current_page_num = slide_idx + 1
            if not (start_page <= current_page_num <= end_page):
                continue

            print(f"--- 正在处理第 {current_page_num} 页 ---")
            
            title = f"未找到标题_第{current_page_num}页"
            text_containing_shapes = [
                s for s in slide.shapes 
                if s.has_text_frame and s.text_frame.text.strip()
            ]
            
            if not text_containing_shapes:
                print(f"警告: 第 {current_page_num} 页没有找到任何文本内容。")
                continue

            sorted_shapes = sorted(text_containing_shapes, key=lambda s: s.top)
            
            title_shape = sorted_shapes[0]
            if title_shape.text_frame.paragraphs:
                title = title_shape.text_frame.paragraphs[0].text.strip()
            else:
                 title = title_shape.text.strip().split('\n')[0]
            
            print(f"识别到标题: {title}")

            slide_content_list = []
            
            for shape in sorted_shapes:
                all_lines = re.split(r'[\n\v\f\r]+', shape.text_frame.text)
                current_key = None
                current_value_lines = []

                def save_previous_item():
                    nonlocal current_key, current_value_lines, slide_content_list
                    if not current_key:
                        return

                    full_value_text = ' '.join(current_value_lines).strip()
                    if not full_value_text:
                        return
                    
                    # 【全新正则表达式】采用“肯定法”定义次级Key
                    # 模式解释:
                    # (?=...)              - 正向先行断言，只匹配位置，不消耗字符。确保分隔符前有标点或空格。
                    #   [。、，,\s]          - 匹配中文句号、顿号、逗号、英文逗号或空格。
                    # (                    - 开始捕获组（这使得分隔符被保留）
                    #   [^\s：:]{2,8}      - Key本身：匹配2到8个非空格、非冒号的字符。
                    #   [：:]              - Key后面的冒号
                    # )                    - 结束捕获组
                    secondary_split_pattern = r'(?=[。、，,\s])\s*([^\s：:]{2,8}[：:])\s*'
                    
                    parts = re.split(secondary_split_pattern, full_value_text)
                    
                    primary_value = parts[0].strip()
                    if primary_value:
                        slide_content_list.append({current_key: primary_value})
                        print(f"  提取到主条目 -> {{'{current_key}': '{primary_value[:70]}...'}}")

                    if len(parts) > 1:
                        # 结果中可能包含空字符串，需要过滤
                        clean_parts = [p for p in parts[1:] if p and p.strip()]
                        for i in range(0, len(clean_parts), 2):
                            secondary_key = clean_parts[i].strip().rstrip('：:')
                            if i + 1 < len(clean_parts):
                                secondary_value = clean_parts[i+1].strip()
                                if secondary_key and secondary_value:
                                    slide_content_list.append({secondary_key: secondary_value})
                                    print(f"  提取到次级条目 -> {{'{secondary_key}': '{secondary_value}'}}")

                    current_key = None
                    current_value_lines = []

                for line in all_lines:
                    clean_line = line.strip()
                    if not clean_line or clean_line == title:
                        continue

                    if clean_line.startswith('■'):
                        save_previous_item()
                        line_to_parse = clean_line.lstrip('■').strip()
                        parts = line_to_parse.replace('：', ':', 1).split(':', 1)
                        if len(parts) == 2:
                            current_key = parts[0].strip()
                            current_value_lines.append(parts[1].strip())
                        else:
                            current_key = line_to_parse
                            current_value_lines = ['']
                    elif current_key:
                        current_value_lines.append(clean_line)
                
                save_previous_item()

            if title in structured_data:
                structured_data[title].extend(slide_content_list)
            else:
                structured_data[title] = slide_content_list

        return structured_data

    except Exception as e:
        import traceback
        print(f"处理PPT文件时发生严重错误: {e}")
        traceback.print_exc()
        return {}

# --- 使用示例 (保持不变) ---
if __name__ == "__main__":
    pptx_file_path = 'E:\project\ppt精细化抽取\关于榆林分公司2025年-2026年员工食堂食材采购项目采购方案的汇报.pptx'
    base_name = os.path.basename(pptx_file_path)
    file_name_without_ext = os.path.splitext(base_name)[0]
    output_json_path = f"2_{file_name_without_ext}_extracted.json"
    start = 3
    end = 5
    
    if not os.path.exists(pptx_file_path):
        print(f"错误：输入文件不存在 -> '{pptx_file_path}'")
    else:
        extracted_data = extract_structured_text_from_pptx(pptx_file_path, start, end)
        print("\n--- 最终提取结果 (控制台预览) ---")
        print(json.dumps(extracted_data, indent=4, ensure_ascii=False))
        if extracted_data:
            try:
                with open(output_json_path, 'w', encoding='utf-8') as f:
                    json.dump(extracted_data, f, indent=4, ensure_ascii=False)
                print(f"\n✅ 提取结果已成功保存到文件: {output_json_path}")
            except IOError as e:
                print(f"\n❌ 保存文件时出错: {e}")
        else:
            print("\n- 未提取到任何数据，不创建输出文件。")