import json
import os
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_structured_text_from_pptx(file_path, start_page=None, end_page=None):
    """
    从具有特定模板的PowerPoint (.pptx)文件中精细化提取结构化文本。
    此版本能正确处理值跨越多行，并使用 '■' 作为新条目的明确分隔符。

    参数:
        file_path (str): PPTX文件的路径。
        start_page (int): 开始提取的页码（从1开始）。
        end_page (int): 结束提取的页码（包含此页）。

    返回:
        dict: 一个字典，键是幻灯片标题，值是一个包含该标题下所有键值对的列表。
    """
    try:
        presentation = Presentation(file_path)
        structured_data = {}

        total_slides = len(presentation.slides)
        if start_page < 1 or end_page > total_slides or start_page > end_page:
            print(f"错误: 页码范围无效。有效范围是 1 到 {total_slides}。")
            return {}
            
        for slide_idx, slide in enumerate(presentation.slides):
            current_page_num = slide_idx + 1
            if not (start_page <= current_page_num <= end_page):
                continue

            print(f"--- 正在处理第 {current_page_num} 页 ---")
            
            title = f"未找到标题_第{current_page_num}页"
            text_containing_shapes = [
                s for s in slide.shapes 
                if s.has_text_frame and s.text_frame.text.strip()
            ]
            
            if not text_containing_shapes:
                print(f"警告: 第 {current_page_num} 页没有找到任何文本内容。")
                continue

            sorted_shapes = sorted(text_containing_shapes, key=lambda s: s.top)
            
            title_shape = sorted_shapes[0]
            if title_shape.text_frame.paragraphs:
                title = title_shape.text_frame.paragraphs[0].text.strip()
            else:
                 title = title_shape.text.strip().split('\n')[0]
            
            print(f"识别到标题: {title}")

            slide_content_list = []
            
            for shape in sorted_shapes:
                # 将整个文本框的内容按行分割
                all_lines = re.split(r'[\n\v\f\r]+', shape.text_frame.text)
                
                current_key = None
                current_value_lines = []

                # 辅助函数，用于处理并保存收集到的一个完整条目
                def save_previous_item():
                    nonlocal current_key, current_value_lines
                    if current_key:
                        # 将多行值合并成一个字符串，用空格连接
                        full_value = ' '.join(current_value_lines).strip()
                        if full_value:
                            slide_content_list.append({current_key: full_value})
                            print(f"  提取到 -> {{'{current_key}': '{full_value[:70]}...'}}") # 打印部分值
                    # 重置状态，为下一个条目做准备
                    current_key = None
                    current_value_lines = []

                for line in all_lines:
                    clean_line = line.strip()
                    if not clean_line or clean_line == title:
                        continue

                    # 【核心逻辑】判断是否为新条目的开始
                    if clean_line.startswith('■'):
                        # 发现新条目标记，说明上一个条目已结束，保存它
                        save_previous_item()

                        # 解析这个新条目
                        line_to_parse = clean_line.lstrip('■').strip()
                        # 统一中英文冒号后进行分割
                        parts = line_to_parse.replace('：', ':', 1).split(':', 1)
                        
                        if len(parts) == 2:
                            current_key = parts[0].strip()
                            current_value_lines.append(parts[1].strip())
                        else:
                            # 处理没有冒号的特殊情况，整个作为key
                            current_key = line_to_parse
                            current_value_lines = ['']
                    elif current_key:
                        # 如果不是新条目的开始，并且我们正在处理一个条目，
                        # 那么这行就是当前值的延续
                        current_value_lines.append(clean_line)
                
                # 循环结束后，不要忘记保存最后一个正在处理的条目
                save_previous_item()

            if title in structured_data:
                structured_data[title].extend(slide_content_list)
            else:
                structured_data[title] = slide_content_list

        return structured_data

    except Exception as e:
        import traceback
        print(f"处理PPT文件时发生严重错误: {e}")
        traceback.print_exc()
        return {}

# --- 使用示例 (保持不变) ---
if __name__ == "__main__":
    # pptx_file_path = 'E:\project\ppt精细化抽取\关于榆林分公司2025年-2026年员工食堂食材采购项目采购方案的汇报.pptx'
    pptx_file_path = 'E:\project\ppt精细化抽取\关于陕西移动2023-2025年风冷型变频氟泵列间机.pptx'
    base_name = os.path.basename(pptx_file_path)
    file_name_without_ext = os.path.splitext(base_name)[0]
    output_json_path = f"{file_name_without_ext}_extracted.json"
    start = 3
    end = 5
    
    if not os.path.exists(pptx_file_path):
        print(f"错误：输入文件不存在 -> '{pptx_file_path}'")
    else:
        extracted_data = extract_structured_text_from_pptx(pptx_file_path, start, end)
        print("\n--- 最终提取结果 (控制台预览) ---")
        print(json.dumps(extracted_data, indent=4, ensure_ascii=False))
        if extracted_data:
            try:
                with open(output_json_path, 'w', encoding='utf-8') as f:
                    json.dump(extracted_data, f, indent=4, ensure_ascii=False)
                print(f"\n✅ 提取结果已成功保存到文件: {output_json_path}")
            except IOError as e:
                print(f"\n❌ 保存文件时出错: {e}")
        else:
            print("\n- 未提取到任何数据，不创建输出文件。")