import json
import os
import re
from pptx import Presentation
# 【新需求】导入 MSO_SHAPE_TYPE 以识别表格和组合形状
from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- 新增辅助函数：专门用于解析表格 ---
def _parse_table_shape(shape):
    """从一个表格形状中提取内容，返回一个二维列表。"""
    table_data = []
    try:
        table = shape.table
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(row_data)
        # 为了与键值对格式统一，我们将表格包装在一个字典中
        if table_data:
            return [{"type": "table", "data": table_data}]
    except Exception as e:
        print(f"  警告: 处理表格时出错: {e}")
    return []

# --- 新增辅助函数：专门用于解析包含 '■' 的文本框 ---
def _parse_key_value_shape(shape):
    """从一个文本框形状中提取基于 '■' 分隔的键值对。"""
    slide_content_list = []
    all_lines = re.split(r'[\n\v\f\r]+', shape.text_frame.text)
    
    current_key = None
    current_value_lines = []

    def save_previous_item():
        nonlocal current_key, current_value_lines
        if current_key:
            # 【优化】处理值内部可能存在的次级键值对
            full_value_text = ' '.join(current_value_lines).strip()
            
            # 使用V7版本的精准正则表达式
            secondary_split_pattern = r'(?=[。、，,\s])\s*([^\s：:]{2,8}[：:])\s*'
            parts = re.split(secondary_split_pattern, full_value_text)
            
            primary_value = parts[0].strip()
            if primary_value:
                slide_content_list.append({current_key: primary_value})
                # print(f"  提取到主条目 -> {{'{current_key}': '{primary_value[:70]}...'}}")

            if len(parts) > 1:
                clean_parts = [p for p in parts[1:] if p and p.strip()]
                for i in range(0, len(clean_parts), 2):
                    secondary_key = clean_parts[i].strip().rstrip('：:')
                    if i + 1 < len(clean_parts):
                        secondary_value = clean_parts[i+1].strip()
                        if secondary_key and secondary_value:
                            slide_content_list.append({secondary_key: secondary_value})
                            # print(f"  提取到次级条目 -> {{'{secondary_key}': '{secondary_value}'}}")
        current_key = None
        current_value_lines = []

    for line in all_lines:
        clean_line = line.strip()
        if not clean_line:
            continue
        if clean_line.startswith('■'):
            save_previous_item()
            line_to_parse = clean_line.lstrip('■').strip()
            parts = line_to_parse.replace('：', ':', 1).split(':', 1)
            if len(parts) == 2:
                current_key = parts[0].strip()
                current_value_lines.append(parts[1].strip())
            else:
                current_key = line_to_parse
                current_value_lines = ['']
        elif current_key:
            current_value_lines.append(clean_line)
    save_previous_item()
    return slide_content_list

# --- 新增辅助函数：处理任意形状（包括递归处理组合形状） ---
def _process_shape(shape):
    """
    递归地处理一个形状，根据其类型分派到相应的解析函数。
    返回一个包含提取内容的列表。
    """
    content = []
    # 情况1: 表格
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        content.extend(_parse_table_shape(shape))
    # 情况2: 组合形状 (需要递归)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            content.extend(_process_shape(sub_shape)) # 递归调用
    # 情况3: 带文本框的形状
    elif shape.has_text_frame and shape.text_frame.text.strip():
        text = shape.text_frame.text
        # 如果文本内容符合您特色的 '■' 格式，则使用精细化解析
        if '■' in text:
            content.extend(_parse_key_value_shape(shape))
        # 否则，作为一个普通文本块提取
        else:
            content.append({"type": "text_block", "content": text.strip()})
    return content

# --- 主函数：架构保持不变，但内容提取逻辑已更新 ---
def extract_structured_text_from_pptx(file_path, start_page=None, end_page=None):
    """
    从PPT文件中提取结构化内容，现在支持键值对、普通文本块和表格。
    """
    try:
        presentation = Presentation(file_path)
        structured_data = {}
        total_slides = len(presentation.slides)

        start = start_page if start_page is not None else 1
        end = end_page if end_page is not None else total_slides
        
        if start < 1 or end > total_slides or start > end:
            print(f"错误: 页码范围无效。有效范围是 1 到 {total_slides}。")
            return {}
            
        for slide_idx, slide in enumerate(presentation.slides):
            current_page_num = slide_idx + 1
            if not (start <= current_page_num <= end):
                continue

            print(f"--- 正在处理第 {current_page_num} 页 ---")
            
            title = f"未找到标题_第{current_page_num}页"
            title_shape_obj = None # 用来存储标题形状对象，以便后续跳过

            # 1. 识别标题 (逻辑保持不变)
            # 找到所有包含文本的形状并按位置排序
            all_text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
            if all_text_shapes:
                sorted_shapes = sorted(all_text_shapes, key=lambda s: s.top)
                title_shape_obj = sorted_shapes[0]
                title = title_shape_obj.text_frame.paragraphs[0].text.strip()
            
            print(f"识别到标题: {title}")

            # 2. 【核心改动】提取页面所有内容（包括表格）
            slide_content_list = []
            # 遍历页面上 *所有* 的形状
            for shape in slide.shapes:
                # 跳过我们已经识别为标题的那个形状
                if shape == title_shape_obj:
                    continue
                
                # 使用新的、强大的处理函数来提取内容
                slide_content_list.extend(_process_shape(shape))

            # 3. 聚合数据 (逻辑保持不变)
            if title in structured_data:
                structured_data[title].extend(slide_content_list)
            else:
                structured_data[title] = slide_content_list

        return structured_data

    except Exception as e:
        import traceback
        print(f"处理PPT文件时发生严重错误: {e}")
        traceback.print_exc()
        return {}

# --- 使用示例 (保持不变) ---
if __name__ == "__main__":
    # 请将路径替换为您的实际文件路径
    # pptx_file_path = 'path/to/your/example.pptx'
    # # pptx_file_path = 'E:\project\ppt精细化抽取\关于榆林分公司2025年-2026年员工食堂食材采购项目采购方案的汇报.pptx'
    pptx_file_path = 'E:\project\优化版本1\ppt精细化抽取\关于咸阳分公司2023-2024年度职工餐厅食材采购.pptx'
   
    base_name = os.path.basename(pptx_file_path)
    file_name_without_ext = os.path.splitext(base_name)[0]
    
    if not os.path.exists(pptx_file_path):
        print(f"错误：输入文件不存在 -> '{pptx_file_path}'")
    else:
        print("\n================== 提取全部页面内容 ==================")
        output_json_path_all = f"{file_name_without_ext}_extracted_all_content.json"
        
        # 调用函数，提取从第3页到结尾的所有内容 (示例)
        extracted_data_all = extract_structured_text_from_pptx(pptx_file_path)
        
        if extracted_data_all:
            print("\n--- 最终提取结果 (全部页面预览) ---")
            print("提取到的标题:", list(extracted_data_all.keys()))
            try:
                with open(output_json_path_all, 'w', encoding='utf-8') as f:
                    json.dump(extracted_data_all, f, indent=4, ensure_ascii=False)
                print(f"\n✅ 全部提取结果已成功保存到文件: {output_json_path_all}")
            except IOError as e:
                print(f"\n❌ 保存文件时出错: {e}")
        else:
            print("\n- 未提取到任何数据，不创建输出文件。")