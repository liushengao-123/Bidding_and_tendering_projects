import json
import os
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- 辅助函数：解析表格 ---
def _parse_table_shape(shape):
    """从表格形状中提取数据，返回二维列表。"""
    table_data = []
    try:
        table = shape.table
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(row_data)
        return table_data
    except Exception as e:
        print(f"  警告: 处理表格时出错: {e}")
    return []

# --- 主函数 ---
def extract_structured_text_from_pptx(file_path, start_page=None, end_page=None):
    """
    V-Final-Corrected: 采用“拆解-重组”策略，根据视觉垂直顺序完美解决内容与Key的关联问题。
    """
    try:
        presentation = Presentation(file_path)
        structured_data = {}
        total_slides = len(presentation.slides)

        start = start_page if start_page is not None else 1
        end = end_page if end_page is not None else total_slides
        
        if start < 1 or end > total_slides or start > end:
            print(f"错误: 页码范围无效。有效范围是 {start}-{end}，但有效范围是 1 到 {total_slides}。")
            return {}
            
        for slide_idx, slide in enumerate(presentation.slides):
            current_page_num = slide_idx + 1
            if not (start <= current_page_num <= end):
                continue

            print(f"--- 正在处理第 {current_page_num} 页 ---")
            
            all_shapes = list(slide.shapes)
            title_shape = None
            if all_shapes:
                sorted_by_top = sorted(all_shapes, key=lambda s: s.top)
                if sorted_by_top and sorted_by_top[0].has_text_frame and sorted_by_top[0].text_frame.text.strip():
                    title_shape = sorted_by_top[0]

            title = title_shape.text_frame.paragraphs[0].text.strip() if title_shape else f"未找到标题_第{current_page_num}页"
            print(f"识别到标题: {title}")

            content_shapes = []
            slide_height = presentation.slide_height
            for shape in all_shapes:
                if shape == title_shape: continue
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if (shape.top / slide_height > 0.9) and text.isdigit() and len(text) <= 3: continue
                content_shapes.append(shape)

            if not content_shapes: continue

            slide_content_list = []
            has_triggers = any(s.has_text_frame and '■' in s.text_frame.text for s in content_shapes)

            if not has_triggers:
                full_page_value = []
                for shape in content_shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table = _parse_table_shape(shape)
                        if table: full_page_value.append({"type": "table", "data": table})
                    elif shape.has_text_frame and shape.text_frame.text.strip():
                        full_page_value.append({"type": "text_block", "content": shape.text_frame.text.strip()})
                if full_page_value: slide_content_list.append({title: full_page_value})
            else:
                # 步骤1：拆解所有内容为带位置的“内容单元”
                content_units = []
                for shape in content_shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_data = _parse_table_shape(shape)
                        if table_data:
                            content_units.append({'top': shape.top, 'type': 'table', 'data': table_data})
                    elif shape.has_text_frame and shape.text_frame.text:
                        text_content = shape.text_frame.text
                        # 如果是触发器文本框，拆分成独立的'■'单元
                        if '■' in text_content:
                            para_count = len(shape.text_frame.paragraphs)
                            shape_height = shape.height if shape.height > 1 else 100 
                            for i, para in enumerate(shape.text_frame.paragraphs):
                                line = para.text.strip()
                                if line.startswith('■'):
                                    estimated_top = shape.top + (i / para_count) * shape_height
                                    content_units.append({'top': estimated_top, 'type': 'key_line', 'text': line})
                        # 否则是普通的文本块
                        elif text_content.strip():
                            content_units.append({'top': shape.top, 'type': 'text_block', 'text': text_content.strip()})
                
                # 步骤2：根据垂直位置对所有“内容单元”进行排序
                sorted_units = sorted(content_units, key=lambda x: x['top'])
                
                # 步骤3：重组
                final_items = []
                for unit in sorted_units:
                    if unit['type'] == 'key_line':
                        line_to_parse = unit['text'].lstrip('■').strip()
                        parts = line_to_parse.replace('：', ':', 1).split(':', 1)
                        key = parts[0].strip()
                        value = parts[1].strip() if len(parts) > 1 else ""
                        final_items.append({key: [value] if value else []})
                    elif (unit['type'] == 'table' or unit['type'] == 'text_block') and final_items:
                        # 将表格或文本块追加到最后一个Key的Value中
                        last_item = final_items[-1]
                        last_item_key = list(last_item.keys())[0]
                        if unit['type'] == 'table':
                            last_item[last_item_key].append({'type': 'table', 'data': unit['data']})
                        else: # text_block
                            last_item[last_item_key].append(unit['text'])
                
                # 步骤4：格式化最终结果
                for item in final_items:
                    key = list(item.keys())[0]
                    value_list = item[key]
                    
                    # 合并连续的字符串
                    merged_list = []
                    temp_text = ""
                    for part in value_list:
                        if isinstance(part, str):
                            temp_text += (" " + part) if temp_text else part
                        else:
                            if temp_text: merged_list.append(temp_text.strip())
                            temp_text = ""
                            merged_list.append(part)
                    if temp_text: merged_list.append(temp_text.strip())

                    # 根据最终列表的长度和内容决定输出格式
                    if len(merged_list) == 1 and isinstance(merged_list[0], str):
                        item[key] = merged_list[0]
                    elif not merged_list:
                        item[key] = ""
                    else:
                        item[key] = merged_list
                slide_content_list = final_items

            if title in structured_data:
                structured_data[title].extend(slide_content_list)
            else:
                structured_data[title] = slide_content_list

        return structured_data
    except Exception as e:
        import traceback
        print(f"处理PPT文件时发生严重错误: {e}")
        traceback.print_exc()
        return {}

# --- 使用示例 ---
if __name__ == "__main__":
    pptx_file_path = 'E:\project\优化版本1\ppt精细化抽取\关于咸阳分公司2023-2024年度职工餐厅食材采购.pptx'
   
    base_name = os.path.basename(pptx_file_path)
    file_name_without_ext = os.path.splitext(base_name)[0]
    
    if not os.path.exists(pptx_file_path):
        print(f"错误：输入文件不存在 -> '{pptx_file_path}'")
    else:
        print("\n================== 提取全部页面内容 ==================")
        output_json_path_all = f"{file_name_without_ext}_extracted_final.json"
        
        extracted_data_all = extract_structured_text_from_pptx(pptx_file_path)
        
        if extracted_data_all:
            print("\n--- 最终提取结果 (预览) ---")
            print("提取到的标题:", list(extracted_data_all.keys()))
            try:
                with open(output_json_path_all, 'w', encoding='utf-8') as f:
                    json.dump(extracted_data_all, f, indent=4, ensure_ascii=False)
                print(f"\n✅ 全部提取结果已成功保存到文件: {output_json_path_all}")
            except IOError as e:
                print(f"\n❌ 保存文件时出错: {e}")
        else:
            print("\n- 未提取到任何数据，不创建输出文件。")