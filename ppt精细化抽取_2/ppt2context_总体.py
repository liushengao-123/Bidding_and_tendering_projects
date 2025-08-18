import json
import os
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- 辅助函数：解析表格 ---
def _parse_table_shape(shape):
    """从表格形状中提取数据，返回二维列表。"""
    table_data = []
    try:
        table = shape.table
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(row_data)
        return table_data
    except Exception as e:
        print(f"  警告: 处理表格时出错: {e}")
    return []

# --- 主函数 ---
def extract_structured_text_from_pptx(file_path, start_page=3, end_page=None):
    """
    V-Final-Fix-5: 采用清晰的状态机逻辑，最终解决所有内容关联问题。
    """
    try:
        presentation = Presentation(file_path)
        structured_data = {}
        total_slides = len(presentation.slides)

        start = start_page if start_page is not None else 1
        end = end_page if end_page is not None else total_slides
        
        if start < 1 or end > total_slides or start > end:
            print(f"错误: 页码范围无效。有效范围是 {start}-{end}，但有效范围是 1 到 {total_slides}。")
            return {}
            
        for slide_idx, slide in enumerate(presentation.slides):
            current_page_num = slide_idx + 1
            if not (start <= current_page_num <= end):
                continue

            print(f"--- 正在处理第 {current_page_num} 页 ---")
            
            all_shapes = sorted(slide.shapes, key=lambda s: (s.top, s.left)) 
            title_shape = None
            if all_shapes and all_shapes[0].has_text_frame and all_shapes[0].text_frame.text.strip():
                title_shape = all_shapes[0]
            title = title_shape.text_frame.paragraphs[0].text.strip() if title_shape else f"未找到标题_第{current_page_num}页"
            print(f"识别到标题: {title}")

            content_shapes = []
            slide_height = presentation.slide_height
            for shape in all_shapes:
                if shape == title_shape: continue
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if (shape.top / slide_height > 0.9) and text.isdigit() and len(text) <= 3: continue
                content_shapes.append(shape)

            if not content_shapes: continue

            slide_content_list = []
            has_triggers = any(s.has_text_frame and '■' in s.text_frame.text for s in content_shapes)

            if not has_triggers:
                full_page_value = []
                for shape in content_shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table = _parse_table_shape(shape)
                        if table: full_page_value.append({"type": "table", "data": table})
                    elif shape.has_text_frame and shape.text_frame.text.strip():
                        full_page_value.append({"type": "text_block", "content": shape.text_frame.text.strip()})
                if full_page_value: slide_content_list.append({title: full_page_value})
            else:
                current_item = None
                
                # 辅助函数，用于保存并格式化一个完整的条目
                def save_item(item):
                    if item:
                        value_parts = item['value_parts']
                        # 合并连续的文本部分
                        final_value_list = []
                        temp_text = ""
                        for part in value_parts:
                            # 确保part是字符串再处理
                            part_text = part if isinstance(part, str) else ""
                            if isinstance(part, str):
                                temp_text += (" " + part_text) if temp_text else part_text
                            else: # 遇到表格等非文本内容
                                if temp_text: 
                                    final_value_list.append(temp_text.strip())
                                    temp_text = ""
                                final_value_list.append(part)
                        if temp_text:
                            final_value_list.append(temp_text.strip())
                        
                        # 根据最终列表的长度和内容决定输出格式
                        final_value = final_value_list[0] if len(final_value_list) == 1 and isinstance(final_value_list[0], str) else final_value_list
                        slide_content_list.append({item['key']: final_value})

                for shape in content_shapes:
                    # 情况1: 当前形状是“触发器”（包含'■'）
                    if shape.has_text_frame and '■' in shape.text_frame.text:
                        lines = re.split(r'[\n\v\f\r]+', shape.text_frame.text)
                        for line in lines:
                            clean_line = line.strip()
                            if clean_line.startswith('■'):
                                # 发现一个新的'■'，意味着上一个条目结束，保存它
                                save_item(current_item)
                                
                                # 开始一个新的条目
                                line_to_parse = clean_line.lstrip('■').strip()
                                parts = line_to_parse.replace('：', ':', 1).split(':', 1)
                                key = parts[0].strip()
                                value = parts[1].strip() if len(parts) > 1 else ""
                                current_item = {'key': key, 'value_parts': [value] if value and value.strip() else []}
                            elif current_item:
                                # 这是当前'■'条目在同一个文本框内的延续文本
                                current_item['value_parts'].append(clean_line)
                    
                    # 情况2: 当前形状不是“触发器”，但我们正在收集一个条目
                    elif current_item:
                        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                            table = _parse_table_shape(shape)
                            if table: current_item['value_parts'].append({"type": "table", "data": table})
                        elif shape.has_text_frame and shape.text_frame.text.strip():
                            current_item['value_parts'].append(shape.text_frame.text.strip())
                    
                    # 情况3: 当前形状不是“触发器”，也没有正在收集的条目（孤立内容）
                    else:
                        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                           table = _parse_table_shape(shape)
                           if table: slide_content_list.append({"type": "table", "data": table})
                        elif shape.has_text_frame and shape.text_frame.text.strip():
                            slide_content_list.append({"type": "text_block", "content": shape.text_frame.text.strip()})
                
                # 循环结束后，不要忘记保存最后一个正在处理的条目
                save_item(current_item)

            if title in structured_data:
                structured_data[title].extend(slide_content_list)
            else:
                structured_data[title] = slide_content_list

        return structured_data
    except Exception as e:
        import traceback
        print(f"处理PPT文件时发生严重错误: {e}")
        traceback.print_exc()
        return {}

# --- 使用示例 ---
if __name__ == "__main__":
    # 请将路径替换为您的实际文件路径
    #pptx_file_path = 'E:\project\优化版本1\ppt精细化抽取\关于咸阳分公司2023-2024年度职工餐厅食材采购.pptx'
    pptx_file_path='E:\project\优化版本1\ppt精细化抽取\关于陕西移动2023-2025年风冷型变频氟泵列间机.pptx'
    base_name = os.path.basename(pptx_file_path)
    file_name_without_ext = os.path.splitext(base_name)[0]
    
    if not os.path.exists(pptx_file_path):
        print(f"错误：输入文件不存在 -> '{pptx_file_path}'")
    else:
        print("\n================== 提取全部页面内容 ==================")
        output_json_path_all = f"E:\project\优化版本1\ppt精细化抽取_2/{file_name_without_ext}_extracted_final.json"
        
        extracted_data_all = extract_structured_text_from_pptx(pptx_file_path)
        
        if extracted_data_all:
            print("\n--- 最终提取结果 (预览) ---")
            print("提取到的标题:", list(extracted_data_all.keys()))
            try:
                with open(output_json_path_all, 'w', encoding='utf-8') as f:
                    json.dump(extracted_data_all, f, indent=4, ensure_ascii=False)
                print(f"\n✅ 全部提取结果已成功保存到文件: {output_json_path_all}")
            except IOError as e:
                print(f"\n❌ 保存文件时出错: {e}")
        else:
            print("\n- 未提取到任何数据，不创建输出文件。")